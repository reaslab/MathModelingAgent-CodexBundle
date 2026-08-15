#!/usr/bin/env python3
# Bundled helper for the document-extraction skill.
import argparse
import os
import sys


MAX_LINE_LENGTH = 2000
MAX_LINE_SUFFIX = f"... (line truncated to {MAX_LINE_LENGTH} chars)"
MAX_BYTES = 50 * 1024
MAX_BYTES_LABEL = f"{MAX_BYTES // 1024} KB"
OFFICE_EXTENSIONS = {
    ".xlsx",
    ".xls",
    ".ods",
    ".docx",
    ".doc",
    ".odt",
    ".pptx",
    ".ppt",
    ".odp",
}


def extract_workbook(filepath: str) -> None:
    from openpyxl import load_workbook

    workbook = load_workbook(filepath, data_only=True, read_only=True)
    for sheet in workbook.worksheets:
        yield f"# Sheet: {sheet.title}"
        for row in sheet.iter_rows(values_only=True):
            yield "\t".join("" if cell is None else str(cell) for cell in row)
        yield ""


def extract_document(filepath: str) -> None:
    from docx import Document

    document = Document(filepath)
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            yield text
    for table in document.tables:
        for row in table.rows:
            yield "\t".join(cell.text for cell in row.cells)


def extract_presentation(filepath: str) -> None:
    from pptx import Presentation

    for index, slide in enumerate(Presentation(filepath).slides, 1):
        yield f"# Slide {index}"
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text.strip():
                yield text
        yield ""


def render_limited(lines, offset: int, limit: int, filetype: str, filepath: str) -> None:
    start = offset - 1
    raw: list[str] = []
    bytes_used = 0
    total_lines = 0
    truncated_by_bytes = False
    has_more_lines = False

    for line in lines:
        total_lines += 1
        if total_lines <= start:
            continue
        if len(raw) >= limit:
            has_more_lines = True
            break

        line = line[:MAX_LINE_LENGTH] + MAX_LINE_SUFFIX if len(line) > MAX_LINE_LENGTH else line
        size = len(line.encode("utf-8")) + (1 if raw else 0)
        if bytes_used + size > MAX_BYTES:
            truncated_by_bytes = True
            has_more_lines = True
            break
        raw.append(line)
        bytes_used += size

    if total_lines < offset and not (total_lines == 0 and offset == 1):
        print(f"Offset {offset} is out of range for extracted {filetype} text ({total_lines} lines)", file=sys.stderr)
        raise SystemExit(4)

    last_read_line = offset + len(raw) - 1
    next_offset = last_read_line + 1
    print(f"<path>{filepath}</path>")
    print(f"<type>office:{filetype}</type>")
    print("<content>")
    for index, line in enumerate(raw, offset):
        print(f"{index}: {line}")
    if truncated_by_bytes:
        print(
            f"\n(Output capped at {MAX_BYTES_LABEL}. Showing extracted lines {offset}-{last_read_line}. Use offset={next_offset} to continue.)"
        )
    elif has_more_lines:
        print(
            f"\n(Showing extracted lines {offset}-{last_read_line}. Use offset={next_offset} to continue.)"
        )
    else:
        print(f"\n(End of extracted content - total {total_lines} lines)")
    print("</content>")


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("filepath")
    parser.add_argument("--offset", type=int, default=1)
    parser.add_argument("--limit", type=int, default=2000)
    args = parser.parse_args()

    filepath = os.path.abspath(args.filepath)
    ext = os.path.splitext(filepath)[1].lower()
    if ext not in OFFICE_EXTENSIONS:
        print(f"Unsupported office format: {ext}", file=sys.stderr)
        return 2

    try:
        if ext in {".xlsx", ".xls", ".ods"}:
            lines = extract_workbook(filepath)
        elif ext in {".docx", ".doc", ".odt"}:
            lines = extract_document(filepath)
        elif ext in {".pptx", ".ppt", ".odp"}:
            lines = extract_presentation(filepath)
        else:
            lines = iter(())
        render_limited(lines, max(1, args.offset), max(1, args.limit), ext.removeprefix("."), filepath)
    except ImportError as exc:
        print(f"missing dependency: {exc}", file=sys.stderr)
        return 3
    except Exception as exc:
        print(f"extract failed: {exc}", file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
